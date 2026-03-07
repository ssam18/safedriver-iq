#!/usr/bin/env python3
"""
Generate SafeDriver-IQ Project Presentation (PPTX)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── Colour palette ──────────────────────────────────────────────
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
ACCENT    = RGBColor(0x00, 0xBF, 0xA5)   # teal accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
BLUE      = RGBColor(0x34, 0x98, 0xDB)
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)
PURPLE    = RGBColor(0x9B, 0x59, 0xB6)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
MEDIUM_BG = RGBColor(0x2C, 0x3E, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper functions ────────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=WHITE, bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(6), space_after=Pt(2),
                  font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_shape_rect(slide, left, top, width, height, fill_color,
                   border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, top=Inches(1.55)):
    """Thin teal accent line under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2.5), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def slide_number_text(n, total=12):
    return f"{n} / {total}"


# ================================================================
#  SLIDE  1  –  TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
            "SafeDriver-IQ", font_size=48, color=ACCENT, bold=True)
add_textbox(slide, Inches(1), Inches(2.4), Inches(10), Inches(0.9),
            "Inverse Crash Modeling for Driver Competency",
            font_size=28, color=WHITE, bold=False)
add_accent_bar(slide, Inches(3.4))

add_textbox(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.6),
            "Quantifying driver safety through continuous scoring (0-100),",
            font_size=20, color=LIGHT_GRAY)
add_textbox(slide, Inches(1), Inches(4.5), Inches(10), Inches(0.6),
            "with special focus on Vulnerable Road User (VRU) protection.",
            font_size=20, color=LIGHT_GRAY)

add_textbox(slide, Inches(1), Inches(5.6), Inches(10), Inches(0.5),
            "Samaresh  |  February 2026  |  github.com/ssam18/safedriver-iq",
            font_size=16, color=LIGHT_GRAY)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(1), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  2  –  THE PROBLEM
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Problem", font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

# Stats boxes
stats = [
    ("7,500+", "Pedestrian deaths\nper year in USA", RED),
    ("1,000+", "Cyclist deaths\nper year in USA", ORANGE),
    ("40-Year", "Record high for\nVRU fatalities", YELLOW),
    ("Reactive", "Current ADAS\nsystems", PURPLE),
]
box_w = Inches(2.6)
box_h = Inches(2.2)
gap = Inches(0.35)
start_x = Inches(0.8)
for i, (num, desc, clr) in enumerate(stats):
    left = start_x + i * (box_w + gap)
    rect = add_shape_rect(slide, left, Inches(1.7), box_w, box_h, MEDIUM_BG, clr)
    add_textbox(slide, left + Inches(0.2), Inches(1.85), box_w - Inches(0.4), Inches(0.8),
                num, font_size=36, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(2.7), box_w - Inches(0.4), Inches(0.9),
                desc, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bullets
bullets = [
    "Traditional systems only warn AFTER danger is detected",
    'No system tells drivers "you\'re safe" or "improve these specific behaviors"',
    "Binary crash predictions (30% risk) aren't actionable for drivers",
]
tf = add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(2.5),
                 "", font_size=18, color=WHITE)
for b in bullets:
    add_paragraph(tf, f"•  {b}", font_size=18, color=LIGHT_GRAY,
                  space_before=Pt(10))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(2), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  3  –  OUR SOLUTION
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Our Solution: Inverse Safety Modeling", font_size=36,
            color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

# Comparison table header
cols = ["", "Traditional Approach", "SafeDriver-IQ (Novel)"]
col_colors = [LIGHT_GRAY, RED, GREEN]
col_x = [Inches(0.8), Inches(3.5), Inches(8.2)]
col_w = [Inches(2.5), Inches(4.5), Inches(4.5)]

rows = [
    ["Metric", "Binary crash prediction", "Continuous safety score (0-100)"],
    ["Output", '"30% crash risk"', '"Safety score: 72/100 → Improve to 85+"'],
    ["Approach", "Reactive warnings", "Proactive guidance with actions"],
    ["Focus", "General risk factors", "VRU-specific safety models"],
]

for j, heading in enumerate(cols):
    clr = ACCENT if j == 0 else WHITE
    add_textbox(slide, col_x[j], Inches(1.5), col_w[j], Inches(0.45),
                heading, font_size=16, color=clr, bold=True)

y_start = Inches(2.1)
for ri, row in enumerate(rows):
    y = y_start + ri * Inches(0.85)
    bg_clr = RGBColor(0x22, 0x2E, 0x3E) if ri % 2 == 0 else MEDIUM_BG
    add_shape_rect(slide, Inches(0.6), y - Inches(0.05), Inches(12), Inches(0.75), bg_clr)
    for j, cell in enumerate(row):
        clr = ACCENT if j == 0 else (RED if j == 1 else GREEN)
        add_textbox(slide, col_x[j], y, col_w[j], Inches(0.65),
                    cell, font_size=16, color=clr, bold=(j == 0))

# Key insight box
rect = add_shape_rect(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.0),
                       RGBColor(0x0A, 0x3D, 0x36), ACCENT)
add_textbox(slide, Inches(1.1), Inches(5.8), Inches(11), Inches(0.8),
            '💡 Key Insight:  Instead of predicting crashes, we compute "distance from crash '
            'boundary" as a continuous safety score.',
            font_size=18, color=ACCENT, bold=True)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(3), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  4  –  DATASET
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Data: CRSS Dataset (2016-2023)", font_size=36,
            color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

data_stats = [
    ("417,335", "Total Crash\nRecords", BLUE),
    ("38,462", "VRU Crashes\n(Ped + Cyclist)", RED),
    ("1,032,571", "Person\nRecords", GREEN),
    ("8 Years", "National NHTSA\nDatabase", ORANGE),
]
for i, (num, desc, clr) in enumerate(data_stats):
    left = start_x + i * (box_w + gap)
    rect = add_shape_rect(slide, left, Inches(1.7), box_w, Inches(2.0), MEDIUM_BG, clr)
    add_textbox(slide, left + Inches(0.2), Inches(1.85), box_w - Inches(0.4), Inches(0.7),
                num, font_size=34, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(2.55), box_w - Inches(0.4), Inches(0.8),
                desc, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Data files
tf = add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.8),
                 "Data Files:", font_size=20, color=ACCENT, bold=True)
files = [
    "ACCIDENT.csv — Crash-level records (~90 columns)",
    "VEHICLE.csv — Vehicle information (~70 columns)",
    "PERSON.csv — Person-level data (~50 columns)",
    "PBTYPE.csv — Pedestrian/bicycle type (~20 columns)",
]
for f in files:
    add_paragraph(tf, f"  •  {f}", font_size=16, color=LIGHT_GRAY, space_before=Pt(8))

# Why this dataset
tf2 = add_textbox(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(2.8),
                  "Why This Dataset?", font_size=20, color=ACCENT, bold=True)
reasons = [
    "National coverage — statistically representative",
    "Rich features — weather, lighting, road, driver info",
    "VRU crash types — dedicated pedestrian/cyclist data",
    "Large scale — 210+ raw columns across all files",
]
for r in reasons:
    add_paragraph(tf2, f"  ✓  {r}", font_size=16, color=LIGHT_GRAY, space_before=Pt(8))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(4), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  5  –  FEATURE ENGINEERING
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Feature Engineering: 210 Raw → 105 Engineered → 57 Selected",
            font_size=32, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

categories = [
    ("Temporal (6)", "HOUR, IS_NIGHT,\nIS_RUSH_HOUR,\nIS_WEEKEND", BLUE),
    ("Environmental (7)", "WEATHER, LGT_COND,\nADVERSE_WEATHER,\nPOOR_LIGHTING", GREEN),
    ("VRU-Specific (7)", "total_vru,\npedestrian_count,\ncyclist_count", RED),
    ("Driver/Vehicle (25)", "SPEED_REL, HARM_EV,\nALCOHOL, MAX_SEV,\nMAN_COLL", ORANGE),
    ("Location (6)", "IS_URBAN,\nHIGH_SPEED_ROAD,\nTYP_INT", PURPLE),
    ("Interactions (3)", "NIGHT_AND_DARK,\nURBAN_HIGH_SPEED,\nWEEKEND_NIGHT", YELLOW),
]

cat_w = Inches(1.85)
cat_h = Inches(2.6)
cat_gap = Inches(0.17)
cat_x0 = Inches(0.55)
for i, (name, examples, clr) in enumerate(categories):
    left = cat_x0 + i * (cat_w + cat_gap)
    rect = add_shape_rect(slide, left, Inches(1.5), cat_w, cat_h, MEDIUM_BG, clr)
    add_textbox(slide, left + Inches(0.1), Inches(1.6), cat_w - Inches(0.2), Inches(0.6),
                name, font_size=14, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.2), cat_w - Inches(0.2), Inches(1.6),
                examples, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Pipeline summary
tf = add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
                 "", font_size=16, color=WHITE)
add_paragraph(tf, "Feature Engineering Pipeline:", font_size=20, color=ACCENT,
              bold=True, space_before=Pt(4))
steps = [
    "~90 raw CRSS columns from ACCIDENT.csv → Added 20 engineered features → 105 total",
    "Removed IDs, non-numeric, target variable, redundant columns → 57 final features",
    "35% of final features are engineered — 5 of top 10 most important are engineered!",
]
for s in steps:
    add_paragraph(tf, f"  →  {s}", font_size=15, color=LIGHT_GRAY, space_before=Pt(6))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(5), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  6  –  HOW INVERSE MODELING WORKS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "How Inverse Modeling Works", font_size=36,
            color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

steps_data = [
    ("1", "Collect Crash Data", "35,314 real VRU crashes\nfrom CRSS database", RED),
    ("2", "Create Safe Samples", "35,314 synthetic safe\nscenarios (modified\ncrash conditions)", GREEN),
    ("3", "Train Classifier", "Binary classification:\nCrash (1) vs Safe (0)\nusing 57 features", BLUE),
    ("4", "Compute Distance", "Distance from decision\nboundary = Safety\nScore (0-100)", ACCENT),
    ("5", "Extract Profile", "\"Good Driver Profile\"\nfrom safest regions\nof feature space", ORANGE),
]

step_w = Inches(2.15)
step_h = Inches(2.8)
step_gap = Inches(0.2)
step_x0 = Inches(0.55)
for i, (num, title, desc, clr) in enumerate(steps_data):
    left = step_x0 + i * (step_w + step_gap)
    rect = add_shape_rect(slide, left, Inches(1.5), step_w, step_h, MEDIUM_BG, clr)
    # Step number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.75), Inches(1.65), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = clr
    circ.line.fill.background()
    add_textbox(slide, left + Inches(0.75), Inches(1.7), Inches(0.6), Inches(0.5),
                num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.4), step_w - Inches(0.2), Inches(0.5),
                title, font_size=14, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.9), step_w - Inches(0.2), Inches(1.2),
                desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Safety score explanation
rect = add_shape_rect(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.2),
                       RGBColor(0x0A, 0x3D, 0x36), ACCENT)
tf = add_textbox(slide, Inches(1.1), Inches(4.8), Inches(11), Inches(2.0),
                 "Safety Score Interpretation:", font_size=18, color=ACCENT, bold=True)
levels = [
    ("0-40  Critical", "Emergency intervention needed", RED),
    ("40-60  High Risk", "Strong warnings issued", ORANGE),
    ("60-75  Medium", "Gentle guidance provided", YELLOW),
    ("75-85  Low Risk", "Minor suggestions only", BLUE),
    ("85-100  Excellent", "Positive feedback given!", GREEN),
]
level_x0 = Inches(1.1)
level_w = Inches(2.1)
for i, (label, desc, clr) in enumerate(levels):
    lx = level_x0 + i * (level_w + Inches(0.1))
    add_textbox(slide, lx, Inches(5.35), level_w, Inches(0.4),
                label, font_size=13, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, lx, Inches(5.7), level_w, Inches(0.5),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(6), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  7  –  THREE MODELS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Model Training: Three Algorithms Compared",
            font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

models = [
    ("Random Forest", "Ensemble of 100 decision trees\ntrained on random data subsets.\n"
     "Predictions averaged for robustness.",
     "n_estimators=100\nmax_depth=10\nmin_samples_split=20",
     "99.38%", "0.9985", BLUE),
    ("XGBoost  ★ WINNER", "Sequential gradient boosting.\nEach tree corrects errors of\nprevious trees. Industry standard.",
     "n_estimators=100\nmax_depth=6\nlearning_rate=0.1",
     "99.45%", "0.9990", GREEN),
    ("Gradient Boosting", "Sklearn's boosting implementation.\nMore conservative, good baseline\nfor comparison.",
     "n_estimators=100\nmax_depth=5\nlearning_rate=0.1",
     "99.30%", "0.9978", PURPLE),
]

model_w = Inches(3.7)
model_h = Inches(4.6)
model_gap = Inches(0.3)
model_x0 = Inches(0.7)
for i, (name, desc, params, acc, auc, clr) in enumerate(models):
    left = model_x0 + i * (model_w + model_gap)
    is_winner = i == 1
    border = ACCENT if is_winner else clr
    rect = add_shape_rect(slide, left, Inches(1.5), model_w, model_h,
                          MEDIUM_BG if not is_winner else RGBColor(0x0A, 0x3D, 0x36), border)
    if is_winner:
        rect.line.width = Pt(3)

    add_textbox(slide, left + Inches(0.2), Inches(1.65), model_w - Inches(0.4), Inches(0.5),
                name, font_size=18, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(2.2), model_w - Inches(0.4), Inches(1.2),
                desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Params
    add_textbox(slide, left + Inches(0.2), Inches(3.3), model_w - Inches(0.4), Inches(0.3),
                "Configuration", font_size=12, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(3.6), model_w - Inches(0.4), Inches(1.0),
                params, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Metrics
    add_textbox(slide, left + Inches(0.2), Inches(4.6), model_w - Inches(0.4), Inches(0.4),
                f"Accuracy: {acc}", font_size=16, color=GREEN if is_winner else WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(5.0), model_w - Inches(0.4), Inches(0.4),
                f"ROC AUC: {auc}", font_size=14, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

# Why train 3?
tf = add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
                 "Why 3 Models?  Different algorithms learn different patterns — comparing "
                 "validates that results aren't algorithm-specific.",
                 font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(7), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  8  –  MODEL RESULTS & FEATURE IMPORTANCE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Results: Feature Importance (XGBoost)",
            font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

# Feature importance as horizontal bars + table
features_imp = [
    ("SPEED_REL", 43.2, "Driver/Vehicle", RED),
    ("total_vru", 12.8, "VRU (Engineered)", GREEN),
    ("ADVERSE_WEATHER", 5.6, "Environmental (Eng)", BLUE),
    ("PEDS", 5.4, "VRU", GREEN),
    ("LGTCON_IM", 5.4, "Environmental", BLUE),
    ("LGT_COND", 4.9, "Environmental", BLUE),
    ("pedestrian_count", 4.8, "VRU (Engineered)", GREEN),
    ("NIGHT_AND_DARK", 3.6, "Interaction (Eng)", YELLOW),
    ("HARM_EV", 2.9, "Driver/Vehicle", ORANGE),
    ("ROAD_COND", 1.9, "Environmental", BLUE),
]

bar_x = Inches(1.0)
bar_max_w = Inches(6.0)
bar_h = Inches(0.38)
bar_y0 = Inches(1.6)
max_val = 43.2

for i, (name, val, cat, clr) in enumerate(features_imp):
    y = bar_y0 + i * (bar_h + Inches(0.08))
    w = bar_max_w * (val / max_val)
    # Label
    add_textbox(slide, Inches(0.1), y, Inches(2.5), bar_h,
                name, font_size=12, color=WHITE, bold=True,
                alignment=PP_ALIGN.RIGHT)
    # Bar
    rect = add_shape_rect(slide, bar_x + Inches(2.0), y, int(w), bar_h, clr)
    # Value
    add_textbox(slide, bar_x + Inches(2.0) + int(w) + Inches(0.1), y,
                Inches(1), bar_h,
                f"{val}%", font_size=12, color=clr, bold=True)

# Right side: insights
insight_x = Inches(8.8)
tf = add_textbox(slide, insight_x, Inches(1.6), Inches(4.0), Inches(5.0),
                 "Key Insights", font_size=20, color=ACCENT, bold=True)
insights = [
    ("Speed dominates (43%)", "Single most important factor.\n"
     "Reflects synthetic safe sample\ncreation process.", RED),
    ("VRU features (23%)", "total_vru + PEDS + pedestrian_\n"
     "count combined. Validates\nour VRU focus.", GREEN),
    ("Environment (16%)", "Weather + lighting combined.\n"
     "Night + darkness interaction\nis particularly important.", BLUE),
    ("Known limitation", "Road condition only 1.9%.\n"
     "Model treats ice = dry.\nFuture: collect real data.", YELLOW),
]
y_ins = Inches(2.2)
for title, desc, clr in insights:
    add_paragraph(tf, f"▸ {title}", font_size=14, color=clr, bold=True,
                  space_before=Pt(12))
    add_paragraph(tf, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(2))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(8), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  9  –  SYSTEM ARCHITECTURE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "System Architecture", font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

# Architecture as layered boxes
layers = [
    ("Data Layer", "CRSS Data (2016-2023)\n417K Crashes → Data Loader → Feature Engineer (57 features)",
     BLUE, Inches(1.5)),
    ("ML Pipeline", "Crash Samples + Synthetic Safe Samples → RF / XGBoost / GB → Best Model Selection",
     GREEN, Inches(2.7)),
    ("Inverse Modeling", "Safety Score Computation (0-100) → Good Driver Profile Extraction",
     ORANGE, Inches(3.9)),
    ("Real-Time System", "Real-Time Calculator → Safety Level → Warnings / Suggestions / Positive Feedback",
     RED, Inches(5.1)),
]

for name, desc, clr, y in layers:
    rect = add_shape_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.0), MEDIUM_BG, clr)
    add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.4),
                name, font_size=16, color=clr, bold=True)
    add_textbox(slide, Inches(3.5), y + Inches(0.05), Inches(8.5), Inches(0.85),
                desc, font_size=14, color=LIGHT_GRAY)

    # Arrow between layers
    if y < Inches(5.0):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(6.5), y + Inches(1.0), Inches(0.4), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Applications bar
rect = add_shape_rect(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7),
                       RGBColor(0x0A, 0x3D, 0x36), ACCENT)
add_textbox(slide, Inches(1.0), Inches(6.35), Inches(11), Inches(0.55),
            "Applications:   Streamlit Dashboard  |  Scenario Simulator  |  SHAP Analysis  |  Agentic AI Agent",
            font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(9), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  10  –  AGENTIC AI
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Agentic AI Integration (Phase 1 Complete)",
            font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "Transforming from passive safety scoring → active autonomous intervention",
            font_size=18, color=LIGHT_GRAY)

agent_modules = [
    ("Perception\nEngine", "Sensor fusion\nVRU detection\nRoad monitoring\n400+ LOC", BLUE),
    ("Decision\nEngine", "6-factor risk\nassessment.\n5 intervention\nlevels. 700+ LOC", GREEN),
    ("Continuous\nLearning", "Experience replay\nOnline RL\nWeight adaptation\n550+ LOC", ORANGE),
    ("Intervention\nController", "Brake control\nMulti-modal alerts\nOverride mgmt\n650+ LOC", RED),
    ("Simulation\nFramework", "6+ scenario types\nRandomized tests\nSequence exec\n300+ LOC", PURPLE),
]

mod_w = Inches(2.15)
mod_h = Inches(2.8)
mod_gap = Inches(0.2)
mod_x0 = Inches(0.55)
for i, (name, desc, clr) in enumerate(agent_modules):
    left = mod_x0 + i * (mod_w + mod_gap)
    rect = add_shape_rect(slide, left, Inches(2.1), mod_w, mod_h, MEDIUM_BG, clr)
    add_textbox(slide, left + Inches(0.1), Inches(2.2), mod_w - Inches(0.2), Inches(0.7),
                name, font_size=15, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(3.0), mod_w - Inches(0.2), Inches(1.6),
                desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: key stats
tf = add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.8),
                 "", font_size=16, color=WHITE)
add_paragraph(tf, "Key Achievement:", font_size=18, color=ACCENT, bold=True)
achievements = [
    "2,650+ lines of production-ready code across 5 modules",
    "Uses existing trained XGBoost model (best_safety_model.pkl) for base safety scores",
    "Multi-factor risk: Safety Score + VRU Proximity + Road + Weather + Driver + Historical",
    "5 intervention levels: Passive → Gentle Warning → Audio Alert → Aggressive → Autonomous Brake",
]
for a in achievements:
    add_paragraph(tf, f"  ✓  {a}", font_size=14, color=LIGHT_GRAY, space_before=Pt(5))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(10), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  11  –  LIMITATIONS & FUTURE WORK
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Limitations & Future Work", font_size=36, color=ACCENT, bold=True)
add_accent_bar(slide, Inches(1.15))

# Left: Limitations
rect = add_shape_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
                       MEDIUM_BG, RED)
tf = add_textbox(slide, Inches(1.1), Inches(1.6), Inches(5.0), Inches(4.8),
                 "Known Limitations", font_size=22, color=RED, bold=True)
limitations = [
    ("Synthetic Safe Data Bias", "Speed dominates at 43% due to how\n"
     "safe samples were generated"),
    ("Road Condition Ignored", "Model treats ice = dry = wet\n"
     "(all show 100% safety)"),
    ("High Accuracy (99.45%)", "Suggests model learned synthetic\n"
     "patterns, not real safety"),
    ("Simulation Only", "No real-world vehicle sensor\n"
     "testing yet"),
]
for title, desc in limitations:
    add_paragraph(tf, f"⚠  {title}", font_size=14, color=ORANGE, bold=True,
                  space_before=Pt(12))
    add_paragraph(tf, f"    {desc}", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

# Right: Future work
rect = add_shape_rect(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.0),
                       MEDIUM_BG, GREEN)
tf = add_textbox(slide, Inches(7.1), Inches(1.6), Inches(5.0), Inches(4.8),
                 "Future Work", font_size=22, color=GREEN, bold=True)
future = [
    ("Short-term (1-2 months)", "Adversarial testing, baseline\n"
     "comparisons, ablation studies"),
    ("Medium-term (3-6 months)", "Collect real driving data (100-500\n"
     "trips), hardware-in-loop testing"),
    ("Long-term (6-12 months)", "Advanced RL (PPO/SAC), pilot\n"
     "deployment, multi-vehicle V2X"),
    ("Publication Target", "IEEE T-ITS (IF: 8.5) or ICRA\n"
     "conference submission"),
]
for title, desc in future:
    add_paragraph(tf, f"→  {title}", font_size=14, color=ACCENT, bold=True,
                  space_before=Pt(12))
    add_paragraph(tf, f"    {desc}", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(11), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ================================================================
#  SLIDE  12  –  THANK YOU / Q&A
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
            "Thank You", font_size=52, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(2.45))

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Questions & Discussion", font_size=28, color=WHITE,
            alignment=PP_ALIGN.CENTER)

# Summary boxes
summary_items = [
    ("417K", "Crash Records", BLUE),
    ("57", "Features", GREEN),
    ("99.45%", "Accuracy", ACCENT),
    ("2,650+", "LOC Agent", ORANGE),
]
sum_w = Inches(2.2)
sum_gap = Inches(0.5)
sum_x0 = Inches(1.5)
for i, (num, label, clr) in enumerate(summary_items):
    left = sum_x0 + i * (sum_w + sum_gap)
    rect = add_shape_rect(slide, left, Inches(4.0), sum_w, Inches(1.3), MEDIUM_BG, clr)
    add_textbox(slide, left, Inches(4.05), sum_w, Inches(0.7),
                num, font_size=28, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(4.7), sum_w, Inches(0.4),
                label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Samaresh  |  github.com/ssam18/safedriver-iq",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
            '"Moving from reactive emergency braking to proactive crash prevention."',
            font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
            slide_number_text(12), font_size=11, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT)


# ── Save ────────────────────────────────────────────────────────
out_path = "SafeDriver-IQ_Presentation.pptx"
prs.save(out_path)
print(f"✅ Presentation saved to {out_path}")
print(f"   Slides: {len(prs.slides)}")
